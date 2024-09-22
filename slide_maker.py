import glob

from pptx import Presentation
from pptx.util import Inches

from effects import (images_appear_after_click_effect,
                     images_appear_on_click_effect)


def create_presentation(filename, images_dir, effect="on_click", delay=None):
    file_extensions = ["jpg", "jpeg", "png"]
    file_extensions += [ext.upper() for ext in file_extensions]  # to include uppercase versions
    images = [filename for ext in file_extensions for filename in glob.glob(f"{images_dir}/*." + ext)]

    top = Inches(1)
    left = Inches(1)
    height = Inches(5.5)

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    for img in images:
        pic = slide.shapes.add_picture(img, left, top, height=height)

    if effect == "on_click":
        # CLICK TO REVEAL EACH IMAGE
        images_appear_on_click_effect(images, slide)

    if effect == "after_click":
        # CLICK ONCE AND EACH IMAGES APPEARS AFTER DELAY (ms)
        images_appear_after_click_effect(images, delay, slide)

    prs.save(filename)


if __name__ == "__main__":
    create_presentation("test.pptx", "nb_images", effect="after_click", delay=200)
